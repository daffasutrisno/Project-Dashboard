"""
PowerPoint presentation builder
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

class PPTBuilder:
    """PowerPoint presentation builder"""
    
    def __init__(self, width=10, height=7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        
        # Layout configuration
        self.chart_width = Inches(2.3)
        self.chart_height = Inches(2.0)
        self.start_left = Inches(0.3)
        self.start_top = Inches(1.0)
        self.h_spacing = Inches(2.4)
        self.v_spacing = Inches(2.1)
    
    def add_slide_with_title(self, title):
        """Add a blank slide with title"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        
        return slide
    
    def add_charts_to_slide(self, slide, charts, chart_list):
        """
        Add charts to slide in 3x4 grid layout
        
        Args:
            slide: PowerPoint slide object
            charts (dict): Dictionary of chart BytesIO streams
            chart_list (list): List of chart names in order
        """
        for idx, chart_name in enumerate(chart_list):
            if chart_name in charts:
                row = idx // 4
                col = idx % 4
                left = self.start_left + (col * self.h_spacing)
                top = self.start_top + (row * self.v_spacing)
                
                slide.shapes.add_picture(
                    charts[chart_name],
                    left, top,
                    width=self.chart_width,
                    height=self.chart_height
                )
    
    def create_5g_slide(self, charts):
        """Create 5G dashboard slide"""
        slide = self.add_slide_with_title('KPI MONITORING 5G EAST JAVA')
        
        chart_list = [
            'availability', 'accessibility', 'cdr', 'sgnb_sr',
            'traffic', 'eut_thp', 'user_5g', 'prb_util',
            'inter_esgnb', 'intra_esgnb', 'intra_sgnb', 'inter_sgnb'
        ]
        
        self.add_charts_to_slide(slide, charts, chart_list)
        return slide
    
    def create_4g_slide(self, charts):
        """Create 4G dashboard slide"""
        slide = self.add_slide_with_title('KPI MONITORING 4G EAST JAVA')
        
        chart_list = [
            'availability', 's1sr', 'rrc_user', 'traffic',
            'eut', 'prb_util', 'cqi', 'qpsk',
            'traffic_split', 'ratio_traffic', 'user_split', 'ratio_user'
        ]
        
        self.add_charts_to_slide(slide, charts, chart_list)
        return slide
    
    def save(self, filename=None):
        """Save presentation to file"""
        if filename is None:
            filename = f'KPI_Monitoring_Dashboard_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
        
        self.prs.save(filename)
        return filename
